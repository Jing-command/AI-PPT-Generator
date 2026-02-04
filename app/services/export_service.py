"""
导出服务
处理 PPTX、PDF、图片导出
"""

import os
import tempfile
import uuid
from pathlib import Path
from typing import Optional

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt

from app.config import settings
from app.models.presentation import Presentation


class ExportService:
    """
    导出服务
    
    支持格式：
    - PPTX: PowerPoint 格式
    - PDF: PDF 格式（需要 LibreOffice 或 Playwright）
    - PNG/JPG: 图片格式
    """
    
    def __init__(self):
        self.storage_path = Path(settings.STORAGE_LOCAL_PATH) / "exports"
        self.storage_path.mkdir(parents=True, exist_ok=True)
    
    async def export_pptx(
        self,
        presentation: Presentation,
        output_path: Optional[str] = None
    ) -> str:
        """
        导出为 PPTX 格式
        
        Args:
            presentation: PPT 模型
            output_path: 输出路径（可选）
            
        Returns:
            导出文件路径
        """
        if output_path is None:
            output_path = str(
                self.storage_path / f"{presentation.id}_{uuid.uuid4().hex}.pptx"
            )
        
        # 创建 PPTX
        prs = PPTXPresentation()
        
        for slide_data in presentation.slides:
            # 根据类型选择布局
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                slide_layout = prs.slide_layouts[0]  # 标题页
            else:
                slide_layout = prs.slide_layouts[1]  # 内容页
            
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            title = slide.shapes.title
            content = slide_data.get('content', {})
            title.text = content.get('title', '')
            
            # 设置内容
            if slide_type != 'title' and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = content.get('text', '')
                
                # 添加要点
                bullets = content.get('bullets', [])
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            
            # 如果有图片，尝试添加
            image_url = content.get('image_url')
            if image_url:
                try:
                    # TODO: 下载图片并添加
                    pass
                except Exception:
                    pass  # 图片添加失败不影响整体导出
        
        # 保存
        prs.save(output_path)
        
        return output_path
    
    async def export_pdf(
        self,
        presentation: Presentation,
        output_path: Optional[str] = None
    ) -> str:
        """
        导出为 PDF 格式
        
        使用 LibreOffice 进行转换
        
        Args:
            presentation: PPT 模型
            output_path: 输出路径
            
        Returns:
            导出文件路径
        """
        # 先生成 PPTX
        pptx_path = await self.export_pptx(presentation)
        
        if output_path is None:
            output_path = str(
                self.storage_path / f"{presentation.id}_{uuid.uuid4().hex}.pdf"
            )
        
        # 使用 LibreOffice 转换
        import subprocess
        
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(Path(output_path).parent),
            pptx_path
        ]
        
        try:
            subprocess.run(cmd, check=True, capture_output=True, timeout=60)
            
            # LibreOffice 生成的文件名可能与预期不同
            generated_pdf = Path(pptx_path).with_suffix('.pdf')
            if generated_pdf.exists():
                generated_pdf.rename(output_path)
            
            # 清理临时 PPTX
            Path(pptx_path).unlink(missing_ok=True)
            
            return output_path
            
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
            # 清理临时文件
            Path(pptx_path).unlink(missing_ok=True)
            raise RuntimeError(f"PDF 导出失败: {e}")
    
    async def export_images(
        self,
        presentation: Presentation,
        format: str = "png",
        output_dir: Optional[str] = None
    ) -> list:
        """
        导出为图片
        
        使用 Playwright 或 LibreOffice 转换
        
        Args:
            presentation: PPT 模型
            format: 图片格式 (png, jpg)
            output_dir: 输出目录
            
        Returns:
            图片路径列表
        """
        # 先生成 PDF
        pdf_path = await self.export_pdf(presentation)
        
        if output_dir is None:
            output_dir = str(
                self.storage_path / f"{presentation.id}_{uuid.uuid4().hex}_images"
            )
        
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)
        
        # 使用 pdf2image 转换
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(
                pdf_path,
                dpi=150,
                fmt=format,
                output_folder=output_dir,
                paths_only=True
            )
            
            # 清理 PDF
            Path(pdf_path).unlink(missing_ok=True)
            
            return images
            
        except ImportError:
            Path(pdf_path).unlink(missing_ok=True)
            raise RuntimeError("pdf2image 未安装，无法导出图片")
    
    def get_file_url(self, file_path: str) -> str:
        """
        获取文件的访问 URL
        
        Args:
            file_path: 文件路径
            
        Returns:
            URL
        """
        # 简化为本地路径，生产环境应使用对象存储
        return f"/exports/{Path(file_path).name}"


# 便捷函数
def get_export_service() -> ExportService:
    """获取导出服务实例"""
    return ExportService()
