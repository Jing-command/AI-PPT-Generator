"""
Export Service
Handle PPTX, PDF, image export
"""

import base64
import os
import re
import tempfile
import uuid
from io import BytesIO
from pathlib import Path
from typing import Optional, Dict, Any
from urllib.parse import urlparse

import httpx
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.config import settings
from app.models.presentation import Presentation


class ImageHelper:
    """Helper class for handling images from URLs or base64 data"""
    
    @staticmethod
    def is_base64(data: str) -> bool:
        """Check if string is base64 encoded image"""
        return data.startswith('data:image') or (len(data) > 100 and not data.startswith('http'))
    
    @staticmethod
    def decode_base64(data: str) -> bytes:
        """Decode base64 image data"""
        if data.startswith('data:image'):
            # Extract base64 part from data URL
            match = re.match(r'data:image/[^;]+;base64,(.+)', data)
            if match:
                data = match.group(1)
        return base64.b64decode(data)
    
    @staticmethod
    async def download_image(url: str) -> bytes:
        """Download image from URL"""
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(url)
            if response.status_code == 200:
                return response.content
            raise ValueError(f"Failed to download image: {response.status_code}")
    
    @staticmethod
    async def get_image_data(image_url: str) -> Optional[BytesIO]:
        """Get image data as BytesIO from URL or base64"""
        try:
            if ImageHelper.is_base64(image_url):
                data = ImageHelper.decode_base64(image_url)
                return BytesIO(data)
            elif image_url.startswith('http'):
                data = await ImageHelper.download_image(image_url)
                return BytesIO(data)
            return None
        except Exception as e:
            print(f"[Export] Failed to get image data: {e}")
            return None


class ExportService:
    """
    Export Service
    
    Supports formats:
    - PPTX: PowerPoint format
    - PDF: PDF format
    - PNG/JPG: Image formats
    """
    
    def __init__(self):
        self.storage_path = Path(settings.STORAGE_LOCAL_PATH) / "exports"
        self.storage_path.mkdir(parents=True, exist_ok=True)
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c*2 for c in hex_color])
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def _add_text_box(self, slide, left, top, width, height, text, 
                      font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        """Helper method to add text box"""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)
        return box
    
    async def export_pptx(
        self,
        presentation: Presentation,
        output_path: Optional[str] = None
    ) -> str:
        """Export to PPTX format"""
        if output_path is None:
            output_path = str(self.storage_path / f"{presentation.id}_{uuid.uuid4().hex}.pptx")
        
        # Create PPTX
        prs = PPTXPresentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for slide_data in presentation.slides:
            slide_type = slide_data.get('type', 'content')
            style = slide_data.get('style', {})
            theme = style.get('theme', {})
            content = slide_data.get('content', {})
            
            # Get theme colors
            primary_rgb = self._hex_to_rgb(theme.get("primary_color", "#1a365d"))
            text_rgb = self._hex_to_rgb(theme.get("text_color", "#1a202c"))
            bg_rgb = self._hex_to_rgb(theme.get("background_color", "#ffffff"))
            
            # Get image URL if available
            image_url = content.get('image_url')
            
            # Create blank slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Set background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*bg_rgb)
            
            # Add background image if available
            image_data = None
            if image_url:
                image_data = await ImageHelper.get_image_data(image_url)
                if image_data:
                    self._add_background_image(slide, image_data)
            
            # Render based on layout type
            if slide_type == 'title':
                self._render_title_slide(slide, content, primary_rgb, text_rgb, image_data is not None)
            elif slide_type == 'section':
                self._render_section_slide(slide, content, primary_rgb, text_rgb, image_data is not None)
            elif slide_type == 'two-column':
                self._render_two_column_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'timeline':
                self._render_timeline_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'process':
                self._render_process_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'grid':
                self._render_grid_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'comparison':
                self._render_comparison_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'data':
                self._render_data_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'quote':
                self._render_quote_slide(slide, content, primary_rgb, text_rgb)
            elif slide_type == 'image-text':
                self._render_image_text_slide(slide, content, primary_rgb, text_rgb, image_data)
            else:
                self._render_content_slide(slide, content, primary_rgb, text_rgb)
        
        prs.save(output_path)
        return output_path
    
    def _add_background_image(self, slide, image_data: BytesIO):
        """Add background image to slide with transparency overlay"""
        try:
            # Add image as background (full slide)
            slide.shapes.add_picture(
                image_data,
                Inches(0), Inches(0),
                width=Inches(13.333),
                height=Inches(7.5)
            )
        except Exception as e:
            print(f"[Export] Failed to add background image: {e}")
    
    def _render_title_slide(self, slide, content, primary_rgb, text_rgb, has_image=False):
        """Render title slide"""
        title = content.get('title', '')
        subtitle = content.get('subtitle', '')
        
        # If has background image, use white text for better visibility
        text_color = RGBColor(255, 255, 255) if has_image else RGBColor(*primary_rgb)
        subtext_color = RGBColor(240, 240, 240) if has_image else RGBColor(*text_rgb)
        
        # Add semi-transparent overlay for better text readability
        if has_image:
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(7.5)
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.4
            overlay.line.fill.background()
        
        self._add_text_box(slide, 0.5, 2.5, 12.333, 1.5, title,
                          font_size=54, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        
        if subtitle:
            self._add_text_box(slide, 0.5, 4.2, 12.333, 1, subtitle,
                              font_size=28, color=subtext_color, align=PP_ALIGN.CENTER)
    
    def _render_section_slide(self, slide, content, primary_rgb, text_rgb, has_image=False):
        """Render section divider slide"""
        title = content.get('title', '')
        description = content.get('description', '')
        
        # If has background image, use white text for better visibility
        text_color = RGBColor(255, 255, 255) if has_image else RGBColor(*primary_rgb)
        subtext_color = RGBColor(240, 240, 240) if has_image else RGBColor(*text_rgb)
        
        # Add semi-transparent overlay for better text readability
        if has_image:
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(7.5)
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.4
            overlay.line.fill.background()
        
        self._add_text_box(slide, 0.5, 2.8, 12.333, 1.5, title,
                          font_size=48, bold=True, color=text_color, align=PP_ALIGN.CENTER)
        
        if description:
            self._add_text_box(slide, 0.5, 4.3, 12.333, 1, description,
                              font_size=24, color=subtext_color, align=PP_ALIGN.CENTER)
    
    def _render_content_slide(self, slide, content, primary_rgb, text_rgb):
        """Render content slide"""
        title = content.get('title', '')
        text = content.get('text', '')
        bullets = content.get('bullets', [])
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if bullets:
            y_pos = 1.5
            for bullet in bullets:
                self._add_text_box(slide, 0.5, y_pos, 12.333, 0.8, f"• {bullet}",
                                  font_size=22, color=text_rgb)
                y_pos += 0.7
        elif text:
            self._add_text_box(slide, 0.5, 1.5, 12.333, 5.5, text,
                              font_size=22, color=text_rgb)
    
    def _render_two_column_slide(self, slide, content, primary_rgb, text_rgb):
        """Render two-column comparison slide"""
        title = content.get('title', '')
        left = content.get('left', {})
        right = content.get('right', {})
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if left.get('title'):
            self._add_text_box(slide, 0.5, 1.4, 5.9, 0.6, left['title'],
                              font_size=28, bold=True, color=primary_rgb)
        
        left_points = left.get('points', [])
        y_pos = 2.2
        for point in left_points[:5]:
            self._add_text_box(slide, 0.5, y_pos, 5.9, 0.6, f"• {point}",
                              font_size=20, color=text_rgb)
            y_pos += 0.6
        
        if right.get('title'):
            self._add_text_box(slide, 6.9, 1.4, 5.9, 0.6, right['title'],
                              font_size=28, bold=True, color=primary_rgb)
        
        right_points = right.get('points', [])
        y_pos = 2.2
        for point in right_points[:5]:
            self._add_text_box(slide, 6.9, y_pos, 5.9, 0.6, f"• {point}",
                              font_size=20, color=text_rgb)
            y_pos += 0.6
    
    def _render_timeline_slide(self, slide, content, primary_rgb, text_rgb):
        """Render timeline slide"""
        title = content.get('title', '')
        events = content.get('events', [])
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if not events:
            return
        
        count = min(len(events), 5)
        spacing = 12.0 / max(count - 1, 1)
        
        for i, event in enumerate(events[:5]):
            x_pos = 0.5 + i * spacing
            
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos - 0.15), Inches(3.85), Inches(0.3), Inches(0.3))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*primary_rgb)
            shape.line.fill.background()
            
            self._add_text_box(slide, x_pos - 0.8, 4.3, 1.6, 0.5,
                              event.get('year', ''), font_size=16, bold=True, 
                              color=primary_rgb, align=PP_ALIGN.CENTER)
            
            self._add_text_box(slide, x_pos - 0.8, 4.8, 1.6, 0.6,
                              event.get('title', ''), font_size=14, bold=True,
                              color=text_rgb, align=PP_ALIGN.CENTER)
            
            self._add_text_box(slide, x_pos - 0.8, 5.3, 1.6, 1.5,
                              event.get('description', ''), font_size=12,
                              color=text_rgb, align=PP_ALIGN.CENTER)
    
    def _render_process_slide(self, slide, content, primary_rgb, text_rgb):
        """Render process flow slide"""
        title = content.get('title', '')
        steps = content.get('steps', [])
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if not steps:
            return
        
        count = min(len(steps), 6)
        spacing = 12.0 / count
        
        for i, step in enumerate(steps[:6]):
            x_pos = 0.5 + i * spacing
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                          Inches(x_pos), Inches(3), Inches(spacing - 0.3), Inches(1.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*primary_rgb)
            shape.line.fill.background()
            
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            
            self._add_text_box(slide, x_pos, 2.3, spacing - 0.3, 0.5,
                              f"Step {i + 1}", font_size=14, color=text_rgb, align=PP_ALIGN.CENTER)
            
            if i < count - 1:
                arrow_x = x_pos + spacing - 0.25
                self._add_text_box(slide, arrow_x, 3.5, 0.2, 0.5, "→",
                                  font_size=24, color=primary_rgb, align=PP_ALIGN.CENTER)
    
    def _render_grid_slide(self, slide, content, primary_rgb, text_rgb):
        """Render grid layout slide"""
        title = content.get('title', '')
        items = content.get('items', [])
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if not items:
            return
        
        positions = [(0.5, 1.5), (6.9, 1.5), (0.5, 4.2), (6.9, 4.2)]
        
        for i, item in enumerate(items[:4]):
            x, y = positions[i]
            
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(x), Inches(y), Inches(5.9), Inches(2.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
            shape.line.color.rgb = RGBColor(*primary_rgb)
            
            self._add_text_box(slide, x + 0.2, y + 0.2, 5.5, 0.6,
                              item.get('title', ''), font_size=24, bold=True, color=primary_rgb)
            
            self._add_text_box(slide, x + 0.2, y + 0.9, 5.5, 1.5,
                              item.get('description', ''), font_size=16, color=text_rgb)
    
    def _render_comparison_slide(self, slide, content, primary_rgb, text_rgb):
        """Render comparison table slide"""
        title = content.get('title', '')
        items = content.get('items', [])
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if not items:
            return
        
        headers = ['Item', 'Option A', 'Option B']
        col_widths = [4, 4.1, 4.1]
        x_pos = 0.5
        
        for i, header in enumerate(headers):
            self._add_text_box(slide, x_pos, 1.5, col_widths[i], 0.6, header,
                              font_size=20, bold=True, color=primary_rgb)
            x_pos += col_widths[i]
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(0.5), Inches(2.1), Inches(12.333), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*primary_rgb)
        line.line.fill.background()
        
        y_pos = 2.3
        for item in items[:6]:
            x_pos = 0.5
            values = [item.get('name', ''), item.get('valueA', ''), item.get('valueB', '')]
            
            for i, value in enumerate(values):
                self._add_text_box(slide, x_pos, y_pos, col_widths[i], 0.5, value,
                                  font_size=18, color=text_rgb)
                x_pos += col_widths[i]
            
            y_pos += 0.6
    
    def _render_data_slide(self, slide, content, primary_rgb, text_rgb):
        """Render data/statistics slide"""
        title = content.get('title', '')
        stats = content.get('stats', [])
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        if not stats:
            return
        
        count = min(len(stats), 4)
        spacing = 12.333 / count
        
        for i, stat in enumerate(stats[:4]):
            x_pos = 0.5 + i * spacing
            
            self._add_text_box(slide, x_pos, 2.5, spacing, 1.5,
                              stat.get('value', ''), font_size=60, bold=True,
                              color=primary_rgb, align=PP_ALIGN.CENTER)
            
            self._add_text_box(slide, x_pos, 4.2, spacing, 0.8,
                              stat.get('label', ''), font_size=20,
                              color=text_rgb, align=PP_ALIGN.CENTER)
    
    def _render_quote_slide(self, slide, content, primary_rgb, text_rgb):
        """Render quote slide"""
        quote = content.get('quote', '')
        author = content.get('author', '')
        title = content.get('title', '')
        
        self._add_text_box(slide, 0.5, 1.5, 1, 1, '"', font_size=120, 
                          color=primary_rgb, align=PP_ALIGN.CENTER)
        
        self._add_text_box(slide, 1.5, 2, 10.333, 2.5, quote,
                          font_size=32, color=text_rgb, align=PP_ALIGN.CENTER)
        
        if author:
            self._add_text_box(slide, 0.5, 4.8, 12.333, 0.8, f"- {author}",
                              font_size=24, color=primary_rgb, align=PP_ALIGN.CENTER)
        
        if title:
            self._add_text_box(slide, 0.5, 5.5, 12.333, 0.6, title,
                              font_size=18, color=text_rgb, align=PP_ALIGN.CENTER)
    
    def _render_image_text_slide(self, slide, content, primary_rgb, text_rgb, image_data=None):
        """Render image-text layout slide"""
        title = content.get('title', '')
        text = content.get('text', '')
        image_url = content.get('image_url', '')
        
        self._add_text_box(slide, 0.5, 0.4, 12.333, 1, title,
                          font_size=40, bold=True, color=primary_rgb)
        
        # Add actual image if available, otherwise use placeholder
        if image_data:
            try:
                slide.shapes.add_picture(
                    image_data,
                    Inches(0.5), Inches(1.5),
                    width=Inches(5.9)
                )
            except Exception as e:
                print(f"[Export] Failed to add image: {e}")
                self._add_image_placeholder(slide, primary_rgb, 0.5, 1.5, 5.9, 5.5)
        else:
            self._add_image_placeholder(slide, primary_rgb, 0.5, 1.5, 5.9, 5.5)
        
        self._add_text_box(slide, 6.9, 1.5, 5.9, 5.5, text,
                          font_size=22, color=text_rgb)
    
    def _add_image_placeholder(self, slide, primary_rgb, left, top, width, height):
        """Add image placeholder shape"""
        img_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(left), Inches(top), Inches(width), Inches(height))
        img_shape.fill.solid()
        img_shape.fill.fore_color.rgb = RGBColor(230, 230, 230)
        img_shape.line.color.rgb = RGBColor(*primary_rgb)
        
        tf = img_shape.text_frame
        p = tf.paragraphs[0]
        p.text = "[Image Placeholder]"
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(150, 150, 150)
    
    async def export_pdf(self, presentation: Presentation, output_path: Optional[str] = None) -> str:
        """Export to PDF"""
        pptx_path = await self.export_pptx(presentation)
        
        if output_path is None:
            output_path = str(self.storage_path / f"{presentation.id}_{uuid.uuid4().hex}.pdf")
        
        import subprocess
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(Path(output_path).parent),
            pptx_path
        ]
        
        try:
            subprocess.run(cmd, check=True, capture_output=True, timeout=60)
            generated_pdf = Path(pptx_path).with_suffix('.pdf')
            if generated_pdf.exists():
                generated_pdf.rename(output_path)
            Path(pptx_path).unlink(missing_ok=True)
            return output_path
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
            Path(pptx_path).unlink(missing_ok=True)
            raise RuntimeError(f"PDF export failed: {e}")
    
    def get_file_url(self, file_path: str) -> str:
        """Get file access URL"""
        return f"/exports/{Path(file_path).name}"


def get_export_service() -> ExportService:
    """Get export service instance"""
    return ExportService()
